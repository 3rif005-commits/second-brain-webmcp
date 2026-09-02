"""Multi-format file extractor.

Supported formats:
  .pdf            → PyMuPDF
  .txt .md .rst   → plain text decode
  .pptx           → python-pptx (slide titles + body text)
  .docx           → python-docx (paragraphs + tables)
"""

import os


def extract_file(path: str, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        from services.pdf_extractor import extract_pdf
        return extract_pdf(path)

    if ext in {".txt", ".md", ".rst", ".csv"}:
        return _extract_text(path)

    if ext == ".pptx":
        return _extract_pptx(path)

    if ext == ".docx":
        return _extract_docx(path)

    raise ValueError(f"Unsupported file type: {ext}")


def _extract_text(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    raise ValueError("Could not decode text file — unknown encoding.")


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text.strip())
                if line.strip():
                    slide_lines.append(line.strip())
        if slide_lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
    if not parts:
        raise ValueError("No readable text found in the presentation.")
    return "\n\n".join(parts)


def _extract_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    lines: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                lines.append(row_text)

    if not lines:
        raise ValueError("No readable text found in the document.")
    return "\n".join(lines)
